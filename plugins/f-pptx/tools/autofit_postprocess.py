#!/usr/bin/env python3
"""PPTX 后处理：给所有文本框加 normAutofit fontScale=55000（文字溢出时自动缩小，不溢出则保持原大）"""
import sys
from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree

FONT_SCALE = sys.argv[2] if len(sys.argv) > 2 else '55000'

prs = Presentation(sys.argv[1])
count = 0

for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        bodyPr = shape.text_frame._txBody.find(qn('a:bodyPr'))
        if bodyPr is None:
            continue
        for tag in [qn('a:normAutofit'), qn('a:spAutoFit'), qn('a:noAutofit')]:
            for el in bodyPr.findall(tag):
                bodyPr.remove(el)
        norm = etree.SubElement(bodyPr, qn('a:normAutofit'))
        norm.set('fontScale', FONT_SCALE)
        count += 1

prs.save(sys.argv[1])
print(f"Done: {count} shapes → normAutofit fontScale={FONT_SCALE}")
